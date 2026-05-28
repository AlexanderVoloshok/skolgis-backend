from __future__ import annotations

import sys
sys.path.insert(1, 'd:/Angular/skolgis-backend')

import os
from pyproj import datadir

proj_dir = datadir.get_data_dir()
os.environ["PROJ_DATA"] = proj_dir
os.environ["PROJ_LIB"] = proj_dir

import matplotlib
matplotlib.use("Agg", force=True)
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
from matplotlib.lines import Line2D
import contextily as cx
import geopandas as gpd
from requests.exceptions import ReadTimeout
from copy import deepcopy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from io import BytesIO
from typing import Tuple, List
from src.config import MAIN_META
from src.aliases import FieldAlias
from src.presentation.utils import get_media_by_fid
from src.sql_utils import read_postgis
from src.utils import timeit

files_table = MAIN_META.tables['skolkovo_general.file_attachments'] 

def iter_shapes_recursive(shapes):
    for sh in shapes:
        yield sh
        # GroupShape имеет shape_type == GROUP и внутри есть .shapes
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(sh.shapes)

def replace_picture(slide, shape_name: str, image_bytes: bytes) -> None:
    """
    Заменяет картинку на слайде на новую, уникальную для этого слайда.
    Сохраняет координаты, размер, имя и alt text.
    """
    for shape in iter_shapes_recursive(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        if shape.name == shape_name or getattr(shape, "alternative_text", "") == shape_name:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            old_name = shape.name
            old_alt = getattr(shape, "alternative_text", "")

            parent = shape.element.getparent()
            old_el = shape.element

            # запомним позицию в дереве, чтобы вставить примерно туда же
            insert_idx = list(parent).index(old_el)

            # удаляем старую картинку
            parent.remove(old_el)

            # вставляем новую
            new_pic = slide.shapes.add_picture(
                BytesIO(image_bytes),
                left,
                top,
                width=width,
                height=height
            )

            # возвращаем имя и alt text
            new_pic.name = old_name
            try:
                new_pic.alternative_text = old_alt or old_name
            except Exception:
                pass

            # перемещаем новую картинку на место старой по z-order
            new_el = new_pic.element
            new_parent = new_el.getparent()
            new_parent.remove(new_el)
            parent.insert(insert_idx, new_el)

            return

    print(f"Picture '{shape_name}' not found on slide")
    return

def _find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    #raise ValueError(f"Shape '{name}' not found on slide")

def fill_attributes(slide, attributes: dict[str, object]) -> None:
    """
    Меняет текст в фигурах по shape.name, сохраняя форматирование.
    Стратегия:
      - если есть runs: пишем весь текст в первый run (стиль сохраняется),
        остальные runs очищаем (чтобы не оставались хвосты старого текста)
      - если runs нет: используем tf.text (это может создать default-стиль,
        но это крайний случай — лучше сделать в шаблоне хотя бы 1 run)
    """
    for shape_name, value in attributes.items():
        shape = _find_shape_by_name(slide, shape_name)
        text = "" if value is None else str(value)

        if not shape:
            continue
        if not shape.has_text_frame:
            print(f"Shape '{shape_name}' has no text_frame")
            continue

        tf = shape.text_frame
        if not tf.paragraphs:
            continue

        p0 = tf.paragraphs[0]

        if p0.runs:
            p0.runs[0].text = text
            for run in p0.runs[1:]:
                run.text = ""

            for p in tf.paragraphs[1:]:
                for run in p.runs:
                    run.text = ""
        else:
            p0.text = text

class PresentationCreator():
    PRESENTATION_TEMPLATE = Presentation("src/assets/slide_sample.pptx")

    def __init__(self, project_ids: List[int]):
        self.project_ids = project_ids
        self.skolkovo_gdf = read_postgis("SELECT * FROM skolkovo_layers.layer_10")

        fields = FieldAlias()
        self.columns_dict = fields.get_field_aliases(orient="dict")

    @timeit
    def render_project_map_png(self, obj_gdf: gpd.GeoDataFrame, surrounding: gpd.GeoDataFrame) -> bytes:
        """
        Рендерит карту:
          - ESRI спутник (World Imagery)
          - границы Сколково красным, без заливки
          - объект (геометрия из obj_gdf)
        Зум: по bounds Сколково.

        Возвращает PNG bytes.
        """

        figsize: Tuple[float, float] = (10, 6)
        dpi: int = 220
        padding_ratio: float = 0.06
        obj_edgecolor: str = "#4240C5"
        obj_lw: float = 1
        obj_facecolor: str = "#5374E1"   # голубой
        surr_edgecolor = "#2F2F2F"
        surr_facecolor = "#696969" 
        sk_borders_color = "#E13838"
        opacity = 0.7

        target_crs = "EPSG:3857"
        sk = self.skolkovo_gdf.to_crs(target_crs)
        obj = obj_gdf.to_crs(target_crs)
        surrounding = surrounding.to_crs(target_crs)

        # bounds по Сколково + паддинг
        minx, miny, maxx, maxy = sk.total_bounds
        dx = maxx - minx
        dy = maxy - miny
        pad_x = dx * padding_ratio
        pad_y = dy * padding_ratio

        fig, ax = plt.subplots(figsize=figsize, dpi=dpi)

        ax.set_xlim(minx - pad_x, maxx + pad_x)
        ax.set_ylim(miny - pad_y, maxy + pad_y)

        # ---------- СПУТНИКОВАЯ OPEN-SOURCE ПОДЛОЖКА ----------
        try:
            cx.add_basemap(ax, source=cx.providers.OpenStreetMap.Mapnik, attribution=False, zorder=0)
        except ReadTimeout:
            pass

        obj.plot(ax=ax, facecolor=obj_facecolor, edgecolor=obj_edgecolor, linewidth=obj_lw, alpha=opacity, zorder=20)
        # Контур поверх — чтобы был насыщенный
        obj.boundary.plot(ax=ax, color=obj_edgecolor, linewidth=obj_lw, alpha=1.0, zorder=25)

        surrounding.plot(ax=ax, facecolor=surr_facecolor, edgecolor=surr_edgecolor, linewidth=obj_lw, alpha=opacity, zorder=20)
        # Контур поверх — чтобы был насыщенный
        surrounding.boundary.plot(ax=ax, color=surr_facecolor, linewidth=obj_lw, alpha=1.0, zorder=25)

        sk.plot(ax=ax, color=sk_borders_color, linewidth=1.3, zorder=23)

        # --- Легенда ---
        # Для проекта легенду делаем как Patch (заливка + контур),
        # для Сколково — как Line2D (только контур).
        handles = [
            Line2D([0], [0], color=sk_borders_color, lw=2.6, label="Границы ИЦ Сколково"),
            Patch(facecolor=obj_facecolor, edgecolor=obj_edgecolor, linewidth=obj_lw, alpha=opacity, label="Границы проекта"),
            Patch(facecolor=surr_facecolor, edgecolor=surr_edgecolor, linewidth=obj_lw, alpha=opacity, label="Окружение"),
        ]
        ax.legend(
            handles=handles, loc="upper left", fontsize=9, framealpha=0.9, borderpad=0.6, labelspacing=0.5, handlelength=2.2
        )

        ax.set_axis_off()
        plt.tight_layout(pad=0)

        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0)
        plt.close(fig)
        buf.seek(0)
        png_bytes = buf.read()

        return png_bytes

    def _prepare_slide_payload(self, project_id: int):
        """
        params: любые параметры для генерации этого слайда
        возвращаем:
          {
            "map_bytes": b"...",
            "attributes": { (row, col): "text", ... }
          }
        """
        obj_gdf = read_postgis(f"""
            SELECT * FROM skolkovo_layers.projects_full
            WHERE project_id = '{project_id}'
        """, crs=4326)
        surrounding = read_postgis(f"""
            WITH target AS (
                SELECT
                    id,
                    project_id,
                    name,
                    year_entered,
                    ST_Transform(ST_MakeValid(geom), 4326) AS g4326
                FROM skolkovo_layers.projects_full
                WHERE project_id = {project_id}
            ),

            nearest AS (
                SELECT DISTINCT ON (p.id)
                    p.id,
                    p.project_id,
                    p.name,
                    p.year_entered,
                    ST_Transform(ST_MakeValid(p.geom), 4326) AS geom,
                    ST_Distance(
                        t.g4326::geography,
                        ST_Transform(ST_MakeValid(p.geom), 4326)::geography
                    ) AS distance_m
                FROM target t
                JOIN skolkovo_layers.projects_full p
                    ON p.id <> t.id
                ORDER BY
                    p.id,
                    ST_Distance(
                        t.g4326::geography,
                        ST_Transform(ST_MakeValid(p.geom), 4326)::geography
                    )
            )

            SELECT
                id,
                project_id,
                name,
                geom,
                year_entered,
                distance_m
            FROM nearest
            ORDER BY distance_m
            LIMIT 2;
        """, crs=4326)
        map_bytes = self.render_project_map_png(obj_gdf, surrounding)
        project_renders = get_media_by_fid(int(obj_gdf.loc[0, 'id']), "render")

        surrRenders = [get_media_by_fid(int(row['project_id']), "render")[0] for _, row in surrounding.iterrows()]
        attrs = obj_gdf.to_dict(orient="records")[0]
        attrs['year_entered'] = int(attrs['year_entered']) if attrs['year_entered'] is not None else None
        attrs['spp_gab'] = int(attrs['spp_gab']) if attrs['spp_gab'] is not None else None
        attrs['parcel_area_ga'] = round(attrs['parcel_area_ga'], 3) if attrs['parcel_area_ga'] else None

        return {
            "map_bytes": map_bytes,
            'project_render_1': project_renders[0][0] if len(project_renders) > 0 else None,
            'project_render_2': project_renders[1][0] if len(project_renders) > 1 else None,
            "SurrRender_1": surrRenders[0][0] if len(surrRenders) > 0 else None,
            "SurrRender_2": surrRenders[1][0] if len(surrRenders) > 1 else None,
            "attributes": attrs,
            "surrounding": surrounding.to_dict(orient="records")
        }
    
    def _clone_slide_template(self):
        """
        Клонирует source_slide в эту же презентацию.
        Внимание: это XML-копирование shapes; подходит для “однотипных” слайдов.

        Ограничения:
          - если на слайде есть диаграммы (charts), media, сложные связи — может потребовать доп. обработки.
          - для картинок/таблиц/текста обычно достаточно.
        """
        source_slide = self.PRESENTATION_TEMPLATE.slides[1]
        layout = source_slide.slide_layout
        new_slide = self.PRESENTATION_TEMPLATE.slides.add_slide(layout)

        # удалить авто-плейсхолдеры от layout
        for shp in list(new_slide.shapes):
            try:
                shp.element.getparent().remove(shp.element)
            except Exception:
                pass

        # копируем shapes по одному
        for shp in source_slide.shapes:
            new_el = deepcopy(shp.element)

            # если это картинка — нужно перенести relationship вручную
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blip = shp._element.blipFill.blip
                old_rid = blip.rEmbed  # в твоей версии именно rEmbed

                rel = source_slide.part.rels[old_rid]

                if rel.is_external:
                    new_rid = new_slide.part.rels._add_relationship(
                        rel.reltype,
                        rel.target_ref,
                        True
                    )
                else:
                    new_rid = new_slide.part.rels._add_relationship(
                        rel.reltype,
                        rel.target_part,
                        False
                    )

                # прописываем новый rId в скопированном XML
                new_blip = new_el.blipFill.blip
                new_blip.set(qn("r:embed"), new_rid)

            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        return new_slide
    
    def _delete_slide(self, slide_index: int) -> None:
        """
        Удаляет слайд по индексу.
        """
        slide_id_list = self.PRESENTATION_TEMPLATE.slides._sldIdLst  # pylint: disable=protected-access
        slide_ids = list(slide_id_list)

        if slide_index < 0 or slide_index >= len(slide_ids):
            raise IndexError("slide_index out of range")

        slide_id = slide_ids[slide_index]
        rId = slide_id.rId

        # remove from slide list
        slide_id_list.remove(slide_id)
        # drop related part
        self.PRESENTATION_TEMPLATE.part.drop_rel(rId)

    def fill_presentation(self):
        payloads = [self._prepare_slide_payload(pid) for pid in self.project_ids]

        # 2) Открываем PPTX и в ОДНОМ потоке обновляем
        for idx, payload in enumerate(payloads):
            slide = self._clone_slide_template()
            # находим картинку и подменяем её blob — без изменения рамки
            replace_picture(slide, "MapImage", payload["map_bytes"])
            fill_attributes(slide, payload['attributes'])

            if payload['project_render_1'] is not None:
                replace_picture(slide, "ProjectRender_1", payload['project_render_1'])
            if payload['project_render_2'] is not None:
                replace_picture(slide, "ProjectRender_2", payload['project_render_2'])

            if payload['SurrRender_1'] is not None:
                replace_picture(slide, "SurrRender_1", payload['SurrRender_1'])
            if payload['SurrRender_2'] is not None:
                replace_picture(slide, "SurrRender_2", payload['SurrRender_2'])

        self._delete_slide(1)
        buf = BytesIO()
        self.PRESENTATION_TEMPLATE.save(buf)
        buf.seek(0)
        return buf
