from __future__ import annotations


import os
from pyproj import datadir

proj_dir = datadir.get_data_dir()
os.environ["PROJ_DATA"] = proj_dir
os.environ["PROJ_LIB"] = proj_dir

import matplotlib
matplotlib.use("Agg", force=True)
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
from matplotlib.lines import Line2D
import contextily as cx
import geopandas as gpd
from copy import deepcopy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from typing import Tuple, List
from src.aliases import FieldAlias
from src.sql_utils import read_postgis
from src.utils import timeit


def iter_shapes_recursive(shapes):
    for sh in shapes:
        yield sh
        # GroupShape имеет shape_type == GROUP и внутри есть .shapes
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(sh.shapes)

def replace_picture(slide, shape_name: str, image_bytes: bytes) -> None:
    for shape in iter_shapes_recursive(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        if shape.name == shape_name or getattr(shape, "alternative_text", "") == shape_name:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height

            # удалить старую фигуру
            shape.element.getparent().remove(shape.element)

            # вставить новую (создаст корректные rels/rId)
            slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=width, height=height)
            return

    raise ValueError(f"Picture '{shape_name}' not found on slide (by name/alt_text).")

def _find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    #raise ValueError(f"Shape '{name}' not found on slide")

def fill_attributes(slide, attributes: dict[str, object]) -> None:
    """
    Меняет текст в фигурах по shape.name, сохраняя форматирование.
    Стратегия:
      - если есть runs: пишем весь текст в первый run (стиль сохраняется),
        остальные runs очищаем (чтобы не оставались хвосты старого текста)
      - если runs нет: используем tf.text (это может создать default-стиль,
        но это крайний случай — лучше сделать в шаблоне хотя бы 1 run)
    """
    for shape_name, value in attributes.items():
        value = round(value, 1) if type(value) == float else value # Все числа округляем до 1 знака
        text = "Нет данных" if value is None else str(value)
        shape = _find_shape_by_name(slide, shape_name)
        if not shape:
            continue

        if not shape.has_text_frame:
            # если вдруг shape не текстовый — пропусти или брось исключение
            raise ValueError(f"Shape '{shape_name}' has no text_frame")

        tf = shape.text_frame
        if not tf.paragraphs:
            tf.text = text
            continue

        # Берём первый параграф
        p0 = tf.paragraphs[0]

        # Если в параграфе уже есть runs — сохраняем стиль первого run
        if p0.runs:
            p0.runs[0].text = text
            # очищаем остальные runs в первом параграфе
            for r in p0.runs[1:]:
                r.text = ""

            # очищаем текст во всех остальных параграфах (если они были)
            for p in tf.paragraphs[1:]:
                for r in p.runs:
                    r.text = ""
        else:
            # В параграфе нет runs: создадим один run (стиль может быть дефолтным)
            run = p0.add_run()
            run.text = text

class PresentationCreator():
    PRESENTATION_TEMPLATE = Presentation("src/assets/slide_sample.pptx")

    def __init__(self, project_ids: List[int]):
        self.project_ids = project_ids
        self.skolkovo_gdf = read_postgis("SELECT * FROM skolkovo_layers.layer_10")

        fields = FieldAlias()
        self.columns_dict = fields.get_field_aliases(orient="dict")

    @timeit
    def render_project_map_png(self, obj_gdf: gpd.GeoDataFrame, surrounding: gpd.GeoDataFrame) -> bytes:
        """
        Рендерит карту:
          - ESRI спутник (World Imagery)
          - границы Сколково красным, без заливки
          - объект (геометрия из obj_gdf)
        Зум: по bounds Сколково.

        Возвращает PNG bytes.
        """

        figsize: Tuple[float, float] = (10, 6)
        dpi: int = 220
        padding_ratio: float = 0.06
        obj_edgecolor: str = "blue"
        obj_lw: float = 2.2
        obj_facecolor: str = "#66CCFF"   # голубой

        target_crs = "EPSG:3857"
        sk = self.skolkovo_gdf.to_crs(target_crs)
        obj = obj_gdf.to_crs(target_crs)
        surrounding = surrounding.to_crs(target_crs)

        # bounds по Сколково + паддинг
        minx, miny, maxx, maxy = sk.total_bounds
        dx = maxx - minx
        dy = maxy - miny
        pad_x = dx * padding_ratio
        pad_y = dy * padding_ratio

        fig, ax = plt.subplots(figsize=figsize, dpi=dpi)

        ax.set_xlim(minx - pad_x, maxx + pad_x)
        ax.set_ylim(miny - pad_y, maxy + pad_y)

        # ---------- СПУТНИКОВАЯ OPEN-SOURCE ПОДЛОЖКА ----------
        cx.add_basemap(ax, source=cx.providers.NASAGIBS.BlueMarble, attribution=False, zorder=0)

        obj.plot(ax=ax, facecolor=obj_facecolor, edgecolor=obj_edgecolor, linewidth=obj_lw, alpha=0.35, zorder=20)
        # Контур поверх — чтобы был насыщенный
        obj.boundary.plot(ax=ax, color=obj_edgecolor, linewidth=obj_lw, alpha=1.0, zorder=25)

        surrounding.plot(ax=ax, facecolor="#9a9a9a", edgecolor="#707070", linewidth=obj_lw, alpha=0.35, zorder=20)
        # Контур поверх — чтобы был насыщенный
        surrounding.boundary.plot(ax=ax, color="#9a9a9a", linewidth=obj_lw, alpha=1.0, zorder=25)

        sk.boundary.plot(ax=ax, color="red", linewidth=2.6, zorder=18)

        # --- Легенда ---
        # Для проекта легенду делаем как Patch (заливка + контур),
        # для Сколково — как Line2D (только контур).
        handles = [
            Line2D([0], [0], color="red", lw=2.6, label="Границы ИЦ Сколково"),
            Patch(facecolor=obj_facecolor, edgecolor=obj_edgecolor, linewidth=obj_lw, alpha=0.35, label="Границы проекта"),
            Patch(facecolor="#9a9a9a", edgecolor="#707070", linewidth=obj_lw, alpha=0.35, label="Окружение"),
        ]
        ax.legend(
            handles=handles, loc="upper left", fontsize=9, framealpha=0.9, borderpad=0.6, labelspacing=0.5, handlelength=2.2
        )

        ax.set_axis_off()
        plt.tight_layout(pad=0)

        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0)
        plt.close(fig)
        buf.seek(0)
        png_bytes = buf.read()

        #out_path = Path('d:/Angular/skolgis-backend/output.png')
        #out_path.parent.mkdir(parents=True, exist_ok=True)
        #out_path.write_bytes(png_bytes)

        return png_bytes

    def _prepare_all_payloads(self):
        """В многопоточном режиме подготавливает данные для каждого слайда.

        Args:
            slide_params_list (_type_): _description_

        Returns:
            _type_: _description_
        """
        #with ThreadPoolExecutor(max_workers=6) as pool:
        #    results = list(pool.map(self._prepare_slide_payload, self.project_ids))
        #return results
        return [self._prepare_slide_payload(pid) for pid in self.project_ids]

    def _prepare_slide_payload(self, project_id: int):
        """
        params: любые параметры для генерации этого слайда
        возвращаем:
          {
            "map_bytes": b"...",
            "attributes": { (row, col): "text", ... }
          }
        """
        obj_gdf = read_postgis(f"""
            SELECT *, round((ST_Area(ST_Transform(ST_MakeValid(geom), 4326)::geography)/ 10000)::numeric,  3) as calc_area 
            FROM skolkovo_layers.projects_full
            WHERE project_id = '{project_id}'
        """, crs=4326)
        surrounding = read_postgis(f"""
            WITH target AS (
              SELECT project_id, name, ST_Transform(ST_MakeValid(geom), 4326) AS g4326 FROM skolkovo_layers.projects_full
              WHERE project_id = '{project_id}'
            )
            SELECT p.project_id, p.name, p.geom, round((ST_Area(ST_Transform(ST_MakeValid(p.geom), 4326)::geography) / 10000)::numeric, 3) AS calc_area,
              -- дистанция в метрах
              round(
                ST_Distance(
                  ST_Transform(ST_MakeValid(p.geom), 4326)::geography,
                  t.g4326::geography
                )::numeric
              , 3) AS distance_m
            FROM skolkovo_layers.projects_full p
            CROSS JOIN target t
            WHERE p.project_id <> t.project_id
            ORDER BY distance_m
            LIMIT 2;
        """, crs=4326)
        map_bytes = self.render_project_map_png(obj_gdf, surrounding)

        #obj_gdf = obj_gdf.rename(columns=self.columns_dict)

        attrs = obj_gdf.to_dict(orient="records")[0]
        #attrs['year_entered'] = int(attrs['year_entered']) if attrs['year_entered'] is not None else None

        return {
            "map_bytes": map_bytes,
            "attributes": attrs,
        }
    
    def _clone_slide_template(self):
        """
        Клонирует source_slide в эту же презентацию.
        Внимание: это XML-копирование shapes; подходит для “однотипных” слайдов.

        Ограничения:
          - если на слайде есть диаграммы (charts), media, сложные связи — может потребовать доп. обработки.
          - для картинок/таблиц/текста обычно достаточно.
        """
        source_slide = self.PRESENTATION_TEMPLATE.slides[1]
        layout = source_slide.slide_layout
        new_slide = self.PRESENTATION_TEMPLATE.slides.add_slide(layout)

        # удалить авто-плейсхолдеры, которые добавил layout
        for shp in list(new_slide.shapes):
            try:
                shp.element.getparent().remove(shp.element)
            except Exception:
                pass

        # копируем shapes из source_slide
        for shp in source_slide.shapes:
            new_el = deepcopy(shp.element)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        return new_slide
    
    def _delete_slide(self, slide_index: int) -> None:
        """
        Удаляет слайд по индексу.
        """
        slide_id_list = self.PRESENTATION_TEMPLATE.slides._sldIdLst  # pylint: disable=protected-access
        slide_ids = list(slide_id_list)

        if slide_index < 0 or slide_index >= len(slide_ids):
            raise IndexError("slide_index out of range")

        slide_id = slide_ids[slide_index]
        rId = slide_id.rId

        # remove from slide list
        slide_id_list.remove(slide_id)
        # drop related part
        self.PRESENTATION_TEMPLATE.part.drop_rel(rId)

    def fill_presentation(self):
        # 1) Параллельно готовим всё тяжёлое
        payloads = self._prepare_all_payloads()

        # 2) Открываем PPTX и в ОДНОМ потоке обновляем
        for idx, payload in enumerate(payloads):
            slide = self._clone_slide_template()
            # находим картинку и подменяем её blob — без изменения рамки
            replace_picture(slide, "MapImage", payload["map_bytes"])
            fill_attributes(slide, payload['attributes'])

        self._delete_slide(1)
        buf = BytesIO()
        self.PRESENTATION_TEMPLATE.save(buf)
        buf.seek(0)
        return buf



if __name__ == '__main__':
    project_ids = [34, 60]

    pres = PresentationCreator(project_ids)
    pres.fill_presentation()

    pres.PRESENTATION_TEMPLATE.save('d:/Angular/skolgis-backend/output.pptx')